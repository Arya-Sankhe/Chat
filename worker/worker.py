import base64
import csv
import json
import os
import random
import re
import shutil
import subprocess
import tempfile
import threading
import time
import uuid
from concurrent.futures import ThreadPoolExecutor, as_completed
from datetime import datetime, timedelta, timezone
from pathlib import Path
from xml.sax.saxutils import escape as xml_escape

import boto3
import edgeparse
import requests
import reportlab
from botocore.config import Config as BotoConfig
from charset_normalizer import from_path
from docx import Document
from docx.shared import Inches as DocxInches
from openpyxl import load_workbook
from openpyxl.utils.cell import get_column_letter, range_boundaries
from pypdf import PdfReader
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptxInches
from pptx.util import Pt
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.platypus import Paragraph, Preformatted, SimpleDocTemplate, Spacer, Table, TableStyle

try:
    from worker.xlsx_generator import create_xlsx_workbook, validate_formula
except ImportError:  # running as a loose script rather than the worker package
    from xlsx_generator import create_xlsx_workbook, validate_formula


def env(name, default=""):
    return os.environ.get(name, default).strip()


def env_int(name, default, minimum=None, maximum=None):
    raw = env(name, str(default))
    try:
        value = int(raw)
    except (TypeError, ValueError):
        value = int(default)
    if minimum is not None:
        value = max(minimum, value)
    if maximum is not None:
        value = min(maximum, value)
    return value


def env_float(name, default, minimum=None, maximum=None):
    raw = env(name, str(default))
    try:
        value = float(raw)
    except (TypeError, ValueError):
        value = float(default)
    if minimum is not None:
        value = max(minimum, value)
    if maximum is not None:
        value = min(maximum, value)
    return value


def default_lease_heartbeat_seconds(lease_seconds):
    lease = max(1, int(lease_seconds or 1))
    return max(5.0, min(float(lease) / 3.0, 30.0))


def is_retryable_http_status(status_code):
    try:
        code = int(status_code)
    except (TypeError, ValueError):
        return False
    return code == 429 or code >= 500


def retry_sleep_seconds(attempt, base=0.5, cap=20.0):
    delay = min(cap, base * (2 ** max(0, attempt)))
    return delay * (0.5 + random.random())


def request_with_retries(method, url, *, max_attempts=3, timeout=30, retry_statuses=True, **kwargs):
    attempts = max(1, int(max_attempts))
    last_error = None
    for attempt in range(attempts):
        try:
            response = requests.request(method, url, timeout=timeout, **kwargs)
            if (
                retry_statuses
                and not response.ok
                and is_retryable_http_status(response.status_code)
                and attempt + 1 < attempts
            ):
                time.sleep(retry_sleep_seconds(attempt))
                continue
            return response
        except requests.exceptions.RequestException as exc:
            last_error = exc
            if attempt + 1 >= attempts:
                raise
            time.sleep(retry_sleep_seconds(attempt))
    if last_error:
        raise last_error
    raise RuntimeError(f"request failed without response: {method} {url}")


def split_page_ranges(page_count, worker_count):
    total = max(0, int(page_count or 0))
    workers = max(1, int(worker_count or 1))
    if total <= 0:
        return []
    if workers <= 1 or total < workers * 2:
        return [(1, total)]
    workers = min(workers, total)
    base, rem = divmod(total, workers)
    ranges = []
    start = 1
    for index in range(workers):
        size = base + (1 if index < rem else 0)
        end = start + size - 1
        ranges.append((start, end))
        start = end + 1
    return ranges


NODE_ARTIFACT_GENERATOR = env("DOCUMENT_ARTIFACT_GENERATOR", str(Path(__file__).resolve().parent / "artifact_generator.mjs"))
NODE_BIN = env("DOCUMENT_NODE_BIN", "node")
USE_JS_ARTIFACT_GENERATOR = env("DOCUMENT_USE_JS_ARTIFACT_GENERATOR", "1").lower() not in ("0", "false", "no")
HTTP_MAX_ATTEMPTS = env_int("DOCUMENT_HTTP_MAX_ATTEMPTS", 3, minimum=1, maximum=8)
WORKER_CONCURRENCY_CAP = 8
PDF_RENDER_WORKERS_CAP = 4
JINA_BATCH_SIZE_CAP = 16
JINA_BATCH_CONCURRENCY_CAP = 4
PAGE_UPLOAD_WORKERS_CAP = 8


class LeaseLostError(RuntimeError):
    pass


class JobCancelledError(RuntimeError):
    def __init__(self, message="job_cancelled"):
        super().__init__(message)
        self.code = "cancelled"


def now_iso():
    return datetime.now(timezone.utc).isoformat()


def normalize_edgeparse_result(raw):
    """Normalize EdgeParse JSON output into a document dict with kids + page count."""
    if raw is None:
        return {"number of pages": 0, "kids": []}

    if isinstance(raw, (bytes, bytearray)):
        raw = raw.decode("utf-8", errors="replace")

    if isinstance(raw, str):
        text = raw.strip()
        if not text:
            return {"number of pages": 0, "kids": []}
        parsed = json.loads(text)
        return normalize_edgeparse_result(parsed)

    if not isinstance(raw, dict):
        for attr in ("json", "data", "document", "result", "output"):
            if hasattr(raw, attr):
                value = getattr(raw, attr)
                if callable(value):
                    try:
                        value = value()
                    except TypeError:
                        pass
                if value is not None and value is not raw:
                    return normalize_edgeparse_result(value)
        if hasattr(raw, "model_dump") and callable(raw.model_dump):
            return normalize_edgeparse_result(raw.model_dump())
        if hasattr(raw, "dict") and callable(raw.dict):
            return normalize_edgeparse_result(raw.dict())
        raise TypeError(f"unsupported EdgeParse result type: {type(raw)!r}")

    data = raw
    if "kids" not in data:
        for key in ("document", "data", "result", "output"):
            nested = data.get(key)
            if isinstance(nested, dict) and ("kids" in nested or "number of pages" in nested):
                data = nested
                break
            if isinstance(nested, str):
                return normalize_edgeparse_result(nested)

    kids = data.get("kids")
    if kids is None:
        kids = []
    if not isinstance(kids, list):
        kids = list(kids)

    page_count = data.get("number of pages")
    if page_count is None:
        page_count = data.get("number_of_pages") or data.get("page_count") or 0
    try:
        page_count = int(page_count or 0)
    except (TypeError, ValueError):
        page_count = 0

    if page_count <= 0 and kids:
        max_page = 0
        for kid in kids:
            if not isinstance(kid, dict):
                continue
            try:
                max_page = max(max_page, int(kid.get("page number") or kid.get("page_number") or 0))
            except (TypeError, ValueError):
                continue
        page_count = max_page

    normalized = dict(data)
    normalized["number of pages"] = page_count
    normalized["kids"] = kids
    return normalized


def _edgeparse_heading_level(element):
    level = element.get("heading level")
    if level is None:
        level = element.get("heading_level")
    if level is None:
        level = element.get("level")
    if isinstance(level, str):
        match = re.search(r"(\d+)", level)
        if match:
            level = match.group(1)
        elif level.strip().lower() == "title":
            level = 1
        else:
            level = 1
    try:
        level = int(level or 1)
    except (TypeError, ValueError):
        level = 1
    return max(1, min(level, 6))


def _edgeparse_cell_text(cell):
    if cell is None:
        return ""
    if isinstance(cell, str):
        return cell.strip()
    if not isinstance(cell, dict):
        return str(cell).strip()
    parts = []
    for kid in cell.get("kids") or []:
        text = _edgeparse_cell_text(kid)
        if text:
            parts.append(text)
    content = cell.get("content")
    if content is not None and str(content).strip():
        parts.append(str(content).strip())
    return " ".join(parts).strip()


def _edgeparse_table_markdown(element):
    rows = element.get("rows") or []
    table_rows = []
    for row in rows:
        if not isinstance(row, dict):
            continue
        cells = row.get("cells") or []
        values = [_edgeparse_cell_text(cell) for cell in cells]
        if any(values):
            table_rows.append(values)
    if not table_rows:
        return ""
    width = max(len(row) for row in table_rows)
    normalized = [normalize_table_row_width(row, width) for row in table_rows]
    header = normalized[0]
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join("---" for _ in header) + " |",
    ]
    for row in normalized[1:]:
        lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines)


def format_edgeparse_element(element):
    if not isinstance(element, dict):
        return ""
    element_type = str(element.get("type") or "").strip().lower()
    content = str(element.get("content") or "").strip()

    if element_type == "heading":
        if not content:
            return ""
        return f"{'#' * _edgeparse_heading_level(element)} {content}"

    if element_type in ("paragraph", "caption", "formula", "text"):
        return content

    if element_type == "list":
        items = []
        for item in element.get("list items") or element.get("list_items") or []:
            if isinstance(item, dict):
                text = str(item.get("content") or "").strip()
            else:
                text = str(item or "").strip()
            if text:
                items.append(f"- {text}")
        return "\n".join(items)

    if element_type == "table":
        return _edgeparse_table_markdown(element)

    if element_type in ("table cell", "table_cell"):
        return _edgeparse_cell_text(element)

    if element_type in ("image", "figure"):
        return ""

    return content


def is_usable_extracted_text(text):
    """Require enough signal that one garbage token cannot unlock text capability."""
    value = (text or "").strip()
    if not value:
        return False
    words = value.split()
    alnum = sum(1 for ch in value if ch.isalnum())
    return len(words) >= 3 or alnum >= 20


def group_edgeparse_pages(document, max_page_count=None):
    """Group EdgeParse kids into page-anchored structured Markdown texts.

    When max_page_count is provided (pypdf page count), it is authoritative:
    elements with page numbers beyond it are ignored and page_count is not inflated.
    """
    doc = normalize_edgeparse_result(document)
    try:
        reported_pages = int(doc.get("number of pages") or 0)
    except (TypeError, ValueError):
        reported_pages = 0
    try:
        max_pages = int(max_page_count) if max_page_count is not None else None
    except (TypeError, ValueError):
        max_pages = None
    if max_pages is not None and max_pages < 0:
        max_pages = 0

    by_page = {}
    for kid in doc.get("kids") or []:
        if not isinstance(kid, dict):
            continue
        try:
            page_number = int(kid.get("page number") or kid.get("page_number") or 0)
        except (TypeError, ValueError):
            continue
        if page_number < 1:
            continue
        if max_pages is not None and page_number > max_pages:
            continue
        by_page.setdefault(page_number, []).append(kid)

    if max_pages is not None:
        page_count = max_pages
    else:
        page_count = reported_pages
        if by_page:
            page_count = max(page_count, max(by_page.keys()))

    pages = []
    for page_number in range(1, page_count + 1):
        elements = by_page.get(page_number) or []
        blocks = []
        for element in elements:
            block = format_edgeparse_element(element)
            if block and block.strip():
                blocks.append(block.strip())
        text = "\n\n".join(blocks).strip()
        pages.append({
            "page_number": page_number,
            "text": text,
            "elements": elements,
            "word_count": len(text.split()) if text else 0,
        })
    combined = "\n\n".join(page["text"] for page in pages if page["text"]).strip()
    return {
        "page_count": page_count,
        "pages": pages,
        "word_count": sum(page["word_count"] for page in pages),
        "has_usable_text": is_usable_extracted_text(combined),
    }


def build_pdftoppm_command(path, prefix, dpi, first=None, last=None):
    cmd = [
        "pdftoppm",
        "-jpeg",
        "-jpegopt",
        "quality=85",
        "-r",
        str(max(72, min(int(dpi or 144), 180))),
    ]
    if first is not None and last is not None:
        cmd.extend(["-f", str(int(first)), "-l", str(int(last))])
    cmd.extend([str(path), str(prefix)])
    return cmd


def lease_until_iso(lease_seconds):
    return (datetime.now(timezone.utc) + timedelta(seconds=max(30, int(lease_seconds or 30)))).isoformat()


def safe_name(value, fallback="document"):
    base = Path(str(value or fallback)).name
    cleaned = "".join(ch if ch.isalnum() or ch in "._-" else "-" for ch in base).strip("-")
    return (cleaned or fallback)[:120]


def truncate(text, limit):
    text = str(text or "")
    return text if len(text) <= limit else text[: max(0, limit - 20)] + "\n...[truncated]"


SUPERSCRIPT_MAP = {
    "⁰": "0", "¹": "1", "²": "2", "³": "3", "⁴": "4",
    "⁵": "5", "⁶": "6", "⁷": "7", "⁸": "8", "⁹": "9",
    "⁺": "+", "⁻": "-", "⁼": "=", "⁽": "(", "⁾": ")",
}
SUBSCRIPT_MAP = {
    "₀": "0", "₁": "1", "₂": "2", "₃": "3", "₄": "4",
    "₅": "5", "₆": "6", "₇": "7", "₈": "8", "₉": "9",
    "₊": "+", "₋": "-", "₌": "=", "₍": "(", "₎": ")",
    "ᵢ": "i", "ⱼ": "j", "ₐ": "a", "ₑ": "e", "ₒ": "o", "ₓ": "x",
}
SYMBOL_TRANSLATION = str.maketrans({
    "ᵀ": "^T",
    "ᵗ": "^t",
    "−": "-",
    "–": "-",
    "—": "-",
    "×": "*",
    "÷": "/",
    "≤": "<=",
    "≥": ">=",
    "≠": "!=",
    "≈": "~=",
    "∈": "in",
    "∉": "not in",
    "∞": "infinity",
})


def normalize_math_symbols(text):
    value = str(text or "").translate(SYMBOL_TRANSLATION)

    def power(match):
        plain = "".join(SUPERSCRIPT_MAP.get(char, char) for char in match.group(0))
        return f"^{plain}" if len(plain) == 1 and plain.isalnum() else f"^({plain})"

    def subscript(match):
        plain = "".join(SUBSCRIPT_MAP.get(char, char) for char in match.group(0))
        return f"_{plain}" if len(plain) == 1 and plain.isalnum() else f"_({plain})"

    value = re.sub(f"[{re.escape(''.join(SUPERSCRIPT_MAP))}]+", power, value)
    value = re.sub(f"[{re.escape(''.join(SUBSCRIPT_MAP))}]+", subscript, value)
    return value


def clean_markdown(text):
    text = str(text or "")
    text = re.sub(r"`([^`]+)`", r"\1", text)
    text = re.sub(r"\*\*([^*]+)\*\*", r"\1", text)
    text = re.sub(r"__([^_]+)__", r"\1", text)
    text = re.sub(r"\*([^*]+)\*", r"\1", text)
    text = re.sub(r"_([^_]+)_", r"\1", text)
    return normalize_math_symbols(text).strip()


def artifact_content(input_data):
    return str(
        input_data.get("content")
        or input_data.get("source_text")
        or input_data.get("data", {}).get("content")
        or input_data.get("data", {}).get("text")
        or input_data.get("data", {}).get("body")
        or ""
    ).strip()


def comparable_heading(text):
    return re.sub(r"[^a-z0-9]+", "", str(text or "").lower())


def strip_duplicate_title_heading(text, title):
    lines = str(text or "").splitlines()
    title_key = comparable_heading(title)
    for index, raw_line in enumerate(lines):
        line = raw_line.strip()
        if not line:
            continue
        heading = re.match(r"^#{1,3}\s+(.+)$", line)
        if heading and comparable_heading(clean_markdown(heading.group(1))) == title_key:
            return "\n".join(lines[:index] + lines[index + 1:]).strip()
        return text
    return text


def split_markdown_table_row(line):
    row = str(line or "").strip()
    if row.startswith("|"):
        row = row[1:]
    if row.endswith("|"):
        row = row[:-1]
    cells = []
    current = []
    escaped = False
    for char in row:
        if escaped:
            current.append(char)
            escaped = False
            continue
        if char == "\\":
            escaped = True
            continue
        if char == "|":
            cells.append("".join(current).strip())
            current = []
            continue
        current.append(char)
    cells.append("".join(current).strip())
    return cells


def is_markdown_table_separator(line):
    cells = split_markdown_table_row(line)
    return len(cells) >= 2 and all(re.match(r"^:?-{3,}:?$", cell.strip()) for cell in cells)


def normalize_table_row_width(row, width):
    values = [str(value or "").strip() for value in row]
    if width <= 0:
        return values
    if len(values) == width:
        return values
    if len(values) < width:
        return values + [""] * (width - len(values))
    if width == 1:
        return ["|".join(values)]
    if width == 2:
        return [values[0], "|".join(values[1:])]
    return values[: width - 2] + ["|".join(values[width - 2:-1])] + [values[-1]]


def collect_markdown_table(lines, start):
    if start + 1 >= len(lines):
        return None, start
    if "|" not in lines[start] or not is_markdown_table_separator(lines[start + 1]):
        return None, start

    headers = split_markdown_table_row(lines[start])
    width = len(headers)
    rows = []
    index = start + 2
    while index < len(lines):
        line = lines[index].strip()
        if not line or "|" not in line:
            break
        rows.append(normalize_table_row_width(split_markdown_table_row(line), width))
        index += 1

    return {"headers": headers, "rows": rows}, index


def find_existing_file(candidates):
    for candidate in candidates:
        path = Path(candidate)
        if path.exists():
            return str(path)
    return None


def register_pdf_fonts():
    reportlab_fonts = Path(reportlab.__file__).resolve().parent / "fonts"
    regular = find_existing_file([
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
        "/Library/Fonts/Arial Unicode.ttf",
        "/System/Library/Fonts/Supplemental/Arial Unicode.ttf",
        Path(__file__).resolve().parent.parent / "fonts" / "DejaVuSans.ttf",
        reportlab_fonts / "Vera.ttf",
    ])
    bold = find_existing_file([
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
        "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf",
        Path(__file__).resolve().parent.parent / "fonts" / "DejaVuSans-Bold.ttf",
        reportlab_fonts / "VeraBd.ttf",
    ])
    mono = find_existing_file([
        "/usr/share/fonts/truetype/dejavu/DejaVuSansMono.ttf",
        "/usr/share/fonts/truetype/liberation/LiberationMono-Regular.ttf",
        Path(__file__).resolve().parent.parent / "fonts" / "DejaVuSansMono.ttf",
        reportlab_fonts / "Vera.ttf",
    ])

    fonts = {"regular": "Helvetica", "bold": "Helvetica-Bold", "mono": "Courier"}
    try:
        if regular:
            pdfmetrics.registerFont(TTFont("KluiSans", regular))
            fonts["regular"] = "KluiSans"
        if bold:
            pdfmetrics.registerFont(TTFont("KluiSans-Bold", bold))
            fonts["bold"] = "KluiSans-Bold"
        elif regular:
            fonts["bold"] = "KluiSans"
        if mono:
            pdfmetrics.registerFont(TTFont("KluiMono", mono))
            fonts["mono"] = "KluiMono"
    except Exception:
        return {"regular": "Helvetica", "bold": "Helvetica-Bold", "mono": "Courier"}
    return fonts


class Supabase:
    def __init__(self):
        self.url = env("SUPABASE_URL").rstrip("/")
        self.key = env("SUPABASE_SERVICE_ROLE_KEY")
        self.max_attempts = HTTP_MAX_ATTEMPTS
        if not self.url or not self.key:
            raise RuntimeError("SUPABASE_URL and SUPABASE_SERVICE_ROLE_KEY are required")

    def _headers(self, prefer=None):
        headers = {
            "apikey": self.key,
            "authorization": f"Bearer {self.key}",
            "content-type": "application/json",
        }
        if prefer:
            headers["prefer"] = prefer
        return headers

    def request(self, path, method="GET", params=None, body=None, prefer=None, retryable=None):
        can_retry = method.upper() in ("GET", "PATCH", "PUT", "DELETE") if retryable is None else bool(retryable)
        response = request_with_retries(
            method,
            f"{self.url}/rest/v1/{path}",
            max_attempts=self.max_attempts if can_retry else 1,
            timeout=30,
            headers=self._headers(prefer),
            params={k: v for k, v in (params or {}).items() if v is not None and v != ""},
            data=json.dumps(body) if body is not None else None,
        )
        if not response.ok:
            raise RuntimeError(f"Supabase {method} {path} failed: {response.status_code} {response.text}")
        if response.status_code == 204 or not response.text:
            return None
        return response.json()

    def rpc(self, name, body):
        return self.request(f"rpc/{name}", method="POST", body=body)

    def claim_job(self, worker_id, lease_seconds):
        rows = self.rpc("klui_claim_document_job", {
            "p_worker_id": worker_id,
            "p_lease_seconds": lease_seconds,
        })
        return rows[0] if rows else None

    def complete_document_job(self, job_id, worker_id, output=None, document_patch=None):
        return self.rpc("klui_complete_document_job", {
            "p_job_id": job_id,
            "p_worker_id": worker_id,
            "p_output": output if output is not None else {},
            "p_document_patch": document_patch if document_patch is not None else {},
        })

    def publish_document_visual_ready(self, job_id, worker_id, document_file_id, page_count, metadata=None):
        return self.rpc("klui_publish_document_visual_ready", {
            "p_job_id": job_id,
            "p_worker_id": worker_id,
            "p_document_file_id": document_file_id,
            "p_page_count": int(page_count),
            "p_metadata": metadata if metadata is not None else {},
        })

    def fail_document_job(self, job_id, worker_id, error):
        return self.rpc("klui_fail_document_job", {
            "p_job_id": job_id,
            "p_worker_id": worker_id,
            "p_error": error if error is not None else {"code": "worker_error", "message": "Document processing failed."},
        })

    def get_job(self, job_id):
        rows = self.request("document_jobs", params={"id": f"eq.{job_id}", "select": "*", "limit": "1"})
        return rows[0] if rows else None

    def get_attachment(self, attachment_id):
        rows = self.request("attachments", params={"id": f"eq.{attachment_id}", "select": "*", "limit": "1"})
        return rows[0] if rows else None

    def get_document_file(self, document_file_id):
        rows = self.request("document_files", params={"id": f"eq.{document_file_id}", "select": "*", "limit": "1"})
        return rows[0] if rows else None

    def update_document_file(self, document_file_id, patch):
        rows = self.request(
            "document_files",
            method="PATCH",
            params={"id": f"eq.{document_file_id}"},
            body={**patch, "updated_at": now_iso()},
            prefer="return=representation",
        )
        return rows[0] if rows else None

    def create_attachment(self, payload):
        rows = self.request("attachments", method="POST", body=payload, prefer="return=representation")
        return rows[0]

    def create_document_file(self, payload):
        rows = self.request("document_files", method="POST", body=payload, prefer="return=representation")
        return rows[0]

    def delete_chunks(self, document_file_id):
        self.request(
            "document_chunks",
            method="DELETE",
            params={"document_file_id": f"eq.{document_file_id}"},
            prefer="return=minimal",
        )

    def delete_pages(self, document_file_id):
        self.request(
            "document_pages",
            method="DELETE",
            params={"document_file_id": f"eq.{document_file_id}"},
            prefer="return=minimal",
        )

    def insert_chunks(self, chunks):
        if not chunks:
            return
        for i in range(0, len(chunks), 250):
            self.request(
                "document_chunks",
                method="POST",
                body=chunks[i:i + 250],
                prefer="resolution=merge-duplicates,return=minimal",
                retryable=True,
            )

    def insert_pages(self, pages, on_conflict="merge"):
        if not pages:
            return
        if on_conflict == "ignore":
            prefer = "resolution=ignore-duplicates,return=minimal"
        else:
            prefer = "resolution=merge-duplicates,return=minimal"
        for i in range(0, len(pages), 100):
            self.request(
                "document_pages",
                method="POST",
                params={"on_conflict": "document_file_id,page_number"},
                body=pages[i:i + 100],
                prefer=prefer,
                retryable=True,
            )

    def update_page(self, document_file_id, page_number, patch):
        rows = self.request(
            "document_pages",
            method="PATCH",
            params={
                "document_file_id": f"eq.{document_file_id}",
                "page_number": f"eq.{int(page_number)}",
            },
            body=patch,
            prefer="return=representation",
        )
        return rows[0] if rows else None

    def update_job(self, job_id, patch, worker_id=None):
        params = {"id": f"eq.{job_id}"}
        if worker_id:
            params["worker_id"] = f"eq.{worker_id}"
        rows = self.request(
            "document_jobs",
            method="PATCH",
            params=params,
            body={**patch, "updated_at": now_iso()},
            prefer="return=representation",
        )
        return rows[0] if rows else None

    def renew_job_lease(self, job_id, worker_id, lease_seconds):
        """Extend lease_until only while the job is still running for this worker."""
        rows = self.request(
            "document_jobs",
            method="PATCH",
            params={
                "id": f"eq.{job_id}",
                "worker_id": f"eq.{worker_id}",
                "status": "eq.running",
                "lease_until": f"gte.{now_iso()}",
            },
            body={
                "lease_until": lease_until_iso(lease_seconds),
                "updated_at": now_iso(),
            },
            prefer="return=representation",
        )
        return rows[0] if rows else None


class JinaEmbeddings:
    def __init__(self):
        self.api_key = env("JINA_API_KEY")
        self.model = env("DOCUMENT_VISUAL_EMBED_MODEL", "jina-embeddings-v5-omni-nano")
        self.dimensions = 768
        self.endpoint = env("JINA_EMBEDDINGS_URL", "https://api.jina.ai/v1/embeddings")
        self.batch_size = env_int("DOCUMENT_JINA_BATCH_SIZE", 8, minimum=1, maximum=JINA_BATCH_SIZE_CAP)
        self.batch_concurrency = env_int(
            "DOCUMENT_JINA_BATCH_CONCURRENCY", 1, minimum=1, maximum=JINA_BATCH_CONCURRENCY_CAP
        )
        self.max_attempts = HTTP_MAX_ATTEMPTS

    @property
    def enabled(self):
        return bool(self.api_key)

    def _embedding_literal(self, values):
        if not values:
            return None
        floats = [float(value) for value in values]
        if len(floats) != 768:
            raise RuntimeError(f"unexpected_embedding_dimensions: got {len(floats)}, expected 768")
        return "[" + ",".join(f"{value:.8g}" for value in floats) + "]"

    def _embed_batch(self, batch):
        inputs = []
        for _, path in batch:
            encoded = base64.b64encode(path.read_bytes()).decode("ascii")
            inputs.append({"bytes": encoded})

        body = {
            "model": self.model,
            "normalized": True,
            "embedding_type": "float",
            "dimensions": self.dimensions,
            "input": inputs,
        }
        response = request_with_retries(
            "POST",
            self.endpoint,
            max_attempts=self.max_attempts,
            timeout=120,
            headers={
                "authorization": f"Bearer {self.api_key}",
                "content-type": "application/json",
            },
            data=json.dumps(body),
        )
        if not response.ok:
            raise RuntimeError(f"Jina embeddings failed: {response.status_code} {response.text[:500]}")

        data = response.json().get("data") or []
        by_index = {int(item.get("index", index)): item.get("embedding") for index, item in enumerate(data)}
        results = []
        for index in range(len(batch)):
            original_index = batch[index][0]
            results.append((original_index, self._embedding_literal(by_index.get(index))))
        return results

    def embed_images(self, image_paths, batch_size=None, batch_concurrency=None):
        if not self.enabled or not image_paths:
            return [None for _ in image_paths]

        size = max(1, min(int(batch_size or self.batch_size), JINA_BATCH_SIZE_CAP))
        concurrency = max(1, min(int(batch_concurrency or self.batch_concurrency), JINA_BATCH_CONCURRENCY_CAP))
        embeddings = [None for _ in image_paths]
        indexed_paths = [
            (index, Path(path))
            for index, path in enumerate(image_paths)
            if Path(path).stat().st_size <= 5 * 1024 * 1024
        ]
        batches = [
            indexed_paths[start:start + size]
            for start in range(0, len(indexed_paths), size)
        ]
        if not batches:
            return embeddings

        if concurrency <= 1 or len(batches) == 1:
            for batch in batches:
                for original_index, literal in self._embed_batch(batch):
                    embeddings[original_index] = literal
            return embeddings

        with ThreadPoolExecutor(max_workers=min(concurrency, len(batches))) as pool:
            futures = [pool.submit(self._embed_batch, batch) for batch in batches]
            for future in as_completed(futures):
                for original_index, literal in future.result():
                    embeddings[original_index] = literal
        return embeddings


class R2:
    def __init__(self):
        account_id = env("R2_ACCOUNT_ID")
        endpoint = env("R2_ENDPOINT") or (f"https://{account_id}.r2.cloudflarestorage.com" if account_id else "")
        self.bucket = env("R2_BUCKET")
        if not endpoint or not self.bucket:
            raise RuntimeError("R2 endpoint and bucket are required")
        self.client = boto3.client(
            "s3",
            endpoint_url=endpoint,
            aws_access_key_id=env("R2_ACCESS_KEY_ID"),
            aws_secret_access_key=env("R2_SECRET_ACCESS_KEY"),
            region_name="auto",
            config=BotoConfig(
                retries={"mode": "adaptive", "max_attempts": max(3, HTTP_MAX_ATTEMPTS)},
            ),
        )

    def download(self, key, path):
        self.client.download_file(self.bucket, key, str(path))

    def upload(self, key, path, content_type):
        with open(path, "rb") as handle:
            response = self.client.put_object(
                Bucket=self.bucket,
                Key=key,
                Body=handle,
                ContentType=content_type,
            )
        return str((response or {}).get("ETag", "")).strip('"')

    def delete(self, key):
        self.client.delete_object(Bucket=self.bucket, Key=key)


class Processor:
    def __init__(self):
        self.db = Supabase()
        self.r2 = R2()
        self.embeddings = JinaEmbeddings()
        self.worker_id = f"document-worker-{uuid.uuid4()}"
        self.lease_seconds = max(30, int(env("DOCUMENT_JOB_TIMEOUT_MS", "120000")) // 1000)
        self.poll_seconds = float(env("DOCUMENT_WORKER_POLL_SECONDS", "1.5"))
        self.max_backoff_seconds = float(env("DOCUMENT_WORKER_MAX_BACKOFF_SECONDS", "30"))
        self.heartbeat_seconds = env_float(
            "DOCUMENT_LEASE_HEARTBEAT_SECONDS",
            default_lease_heartbeat_seconds(self.lease_seconds),
            minimum=1.0,
        )
        self.heartbeat_seconds = min(self.heartbeat_seconds, max(1.0, self.lease_seconds / 2.0))
        self._lease_lost = None
        self._active_job_id = None
        self.max_extracted_chars = int(env("DOCUMENT_MAX_EXTRACTED_CHARS", "500000"))
        self.visual_page_dpi = int(env("DOCUMENT_VISUAL_PAGE_DPI", "144"))
        self.pdf_render_workers = env_int(
            "DOCUMENT_PDF_RENDER_WORKERS", 2, minimum=1, maximum=PDF_RENDER_WORKERS_CAP
        )
        self.page_upload_workers = env_int(
            "DOCUMENT_PAGE_UPLOAD_WORKERS", 4, minimum=1, maximum=PAGE_UPLOAD_WORKERS_CAP
        )
        self.default_limits = {
            "max_pdf_pages": int(env("DOCUMENT_MAX_PDF_PAGES", "100")),
            "max_docx_words": int(env("DOCUMENT_MAX_DOCX_WORDS", "80000")),
            "max_xlsx_sheets": int(env("DOCUMENT_MAX_XLSX_SHEETS", "25")),
            "max_xlsx_cells": int(env("DOCUMENT_MAX_XLSX_CELLS", "250000")),
            "max_csv_rows": int(env("DOCUMENT_MAX_CSV_ROWS", "100000")),
            "max_csv_columns": int(env("DOCUMENT_MAX_CSV_COLUMNS", "100")),
            "max_extracted_chars": self.max_extracted_chars,
        }

    def object_key(self, user_id, file_name):
        return f"users/{user_id}/{uuid.uuid4()}/{safe_name(file_name)}"

    def run(self):
        print(f"{self.worker_id} started", flush=True)
        claim_failures = 0
        while True:
            try:
                job = self.db.claim_job(self.worker_id, self.lease_seconds)
                claim_failures = 0
                if not job:
                    time.sleep(self.poll_seconds)
                    continue
                self.handle_job(job)
            except requests.exceptions.RequestException as exc:
                claim_failures += 1
                sleep_for = min(
                    self.max_backoff_seconds,
                    self.poll_seconds * (2 ** min(claim_failures - 1, 5)),
                )
                print(f"worker claim network error; retrying in {sleep_for:.1f}s: {exc}", flush=True)
                time.sleep(sleep_for)
            except Exception as exc:
                print(f"worker loop error: {exc}", flush=True)
                time.sleep(self.poll_seconds)

    def _lease_heartbeat_loop(self, job_id, stop_event, lost_event):
        last_renewed = time.monotonic()
        while not stop_event.wait(self.heartbeat_seconds):
            try:
                renewed = self.db.renew_job_lease(job_id, self.worker_id, self.lease_seconds)
                if not renewed:
                    print(f"job {job_id} lease heartbeat skipped (not running for {self.worker_id})", flush=True)
                    lost_event.set()
                    return
                last_renewed = time.monotonic()
            except Exception as exc:
                print(f"job {job_id} lease heartbeat failed: {exc}", flush=True)
                safe_window = max(1.0, self.lease_seconds - self.heartbeat_seconds)
                if time.monotonic() - last_renewed >= safe_window:
                    print(f"job {job_id} stopped before its unrenewed lease could expire", flush=True)
                    lost_event.set()
                    return

    def handle_job(self, job):
        job_id = job["id"]
        tmp = Path(tempfile.mkdtemp(prefix=f"doc-job-{job_id}-"))
        stop_heartbeat = threading.Event()
        lease_lost = threading.Event()
        self._lease_lost = lease_lost
        self._active_job_id = job_id
        heartbeat = threading.Thread(
            target=self._lease_heartbeat_loop,
            args=(job_id, stop_heartbeat, lease_lost),
            name=f"lease-heartbeat-{job_id}",
            daemon=True,
        )
        heartbeat.start()
        try:
            output = self.dispatch(job, tmp)
            self.assert_job_active(job_id)
            document_patch = {}
            if isinstance(output, dict):
                document_patch = output.pop("_document_patch", None) or {}
                public_output = output
            else:
                public_output = output or {}
            completed = self.db.complete_document_job(
                job_id,
                self.worker_id,
                public_output,
                document_patch,
            )
            if completed is None:
                raise LeaseLostError("job_lease_lost")
        except LeaseLostError as exc:
            print(f"job {job_id} stopped after losing its lease: {exc}", flush=True)
        except JobCancelledError as exc:
            error = {"message": str(exc), "code": getattr(exc, "code", "cancelled")}
            failed = self.db.fail_document_job(job_id, self.worker_id, error)
            if failed is None:
                print(f"job {job_id} cancel ignored after losing its lease", flush=True)
            else:
                print(f"job {job_id} cancelled: {exc}", flush=True)
        except Exception as exc:
            error = {"message": str(exc), "code": getattr(exc, "code", "worker_error")}
            failed = self.db.fail_document_job(job_id, self.worker_id, error)
            if failed is None:
                print(f"job {job_id} failure ignored after losing its lease: {exc}", flush=True)
            else:
                print(f"job {job_id} failed: {exc}", flush=True)
        finally:
            stop_heartbeat.set()
            heartbeat.join(timeout=max(1.0, min(self.heartbeat_seconds, 5.0)))
            self._lease_lost = None
            self._active_job_id = None
            shutil.rmtree(tmp, ignore_errors=True)

    def assert_job_lease(self):
        if self._lease_lost is not None and self._lease_lost.is_set():
            raise LeaseLostError("job_lease_lost")

    def assert_job_active(self, job_id=None):
        self.assert_job_lease()
        target_id = job_id or getattr(self, "_active_job_id", None)
        if not target_id:
            return
        current = self.db.get_job(target_id)
        if not current:
            raise LeaseLostError("job_missing")
        if current.get("cancel_requested"):
            raise JobCancelledError("job_cancelled")
        if current.get("status") != "running" or current.get("worker_id") != self.worker_id:
            raise LeaseLostError("job_lease_lost")

    def dispatch(self, job, tmp):
        job_type = job["job_type"]
        if job_type == "document.enrich.pdf":
            return self.enrich_pdf_job(job, tmp)
        if job_type == "document.render_page":
            return self.render_page_job(job, tmp)
        if job_type.startswith("document.extract."):
            return self.extract_job(job, tmp)
        if job_type.startswith("document.create."):
            return self.create_job(job, tmp)
        if job_type.startswith("document.edit."):
            return self.edit_job(job, tmp)
        if job_type.startswith("document.export."):
            return self.export_job(job, tmp)
        raise RuntimeError(f"Unsupported job type: {job_type}")

    def extract_job(self, job, tmp):
        doc = self.db.get_document_file(job["document_file_id"])
        attachment = self.db.get_attachment(doc["attachment_id"])
        source = tmp / safe_name(attachment["file_name"])
        self.r2.download(attachment["object_key"], source)
        self.assert_job_active(job["id"])
        limits = {**self.default_limits, **((job.get("input") or {}).get("limits") or {})}
        if doc["kind"] == "pdf":
            return self.extract_pdf_text_job(job, tmp, doc, attachment, source, limits)

        chunks, meta = self.extract(source, doc["kind"], doc["user_id"], doc["id"], limits)
        extraction = {
            "metadata": meta,
            "chunks": [{k: v for k, v in chunk.items() if k not in ("user_id", "document_file_id")} for chunk in chunks],
        }
        extraction_path = tmp / "extraction.json"
        extraction_path.write_text(json.dumps(extraction, ensure_ascii=False), encoding="utf-8")
        extraction_key = self.object_key(doc["user_id"], f"{Path(attachment['file_name']).stem}.extraction.json")
        self.r2.upload(extraction_key, extraction_path, "application/json")
        self.assert_job_active(job["id"])
        self.db.delete_chunks(doc["id"])
        self.db.insert_chunks(chunks)
        self.assert_job_active(job["id"])
        ready_at = now_iso()
        meta_out = {
            **(doc.get("metadata") or {}),
            **meta,
            "progress": 100,
            "stage": "text_ready",
            "file_name": attachment["file_name"],
            "content_type": attachment["content_type"],
            "size_bytes": attachment["size_bytes"],
        }
        return {
            "document_file_id": doc["id"],
            "status": "text_ready",
            **meta,
            "_document_patch": {
                "text_ready_at": ready_at,
                "word_count": meta.get("word_count"),
                "extraction_key": extraction_key,
                "metadata": meta_out,
                "error": None,
            },
        }

    def extract_pdf_text_job(self, job, tmp, doc, attachment, source, limits):
        job_started = time.monotonic()
        self.db.update_document_file(doc["id"], {
            "processing_status": "processing",
            "metadata": {
                "mode": "edgeparse_text",
                "progress": 5,
                "stage": "extracting_text",
                "file_name": attachment["file_name"],
                "content_type": attachment["content_type"],
                "size_bytes": attachment["size_bytes"],
            },
        })
        self.assert_job_active(job["id"])

        reader = PdfReader(str(source))
        if reader.is_encrypted:
            raise RuntimeError("password_protected")
        page_count = len(reader.pages)
        max_pages = self.limit(limits, "max_pdf_pages", 100)
        if page_count > max_pages:
            raise RuntimeError(f"too_many_pages: PDF has {page_count} pages; limit is {max_pages}")

        extract_started = time.monotonic()
        chunks, meta = self.extract_pdf(source, doc["user_id"], doc["id"], limits, page_count=page_count)
        extract_seconds = time.monotonic() - extract_started
        self.assert_job_active(job["id"])

        extraction = {
            "metadata": meta,
            "chunks": [{k: v for k, v in chunk.items() if k not in ("user_id", "document_file_id")} for chunk in chunks],
        }
        extraction_path = tmp / "extraction.json"
        extraction_path.write_text(json.dumps(extraction, ensure_ascii=False), encoding="utf-8")
        extraction_key = self.object_key(doc["user_id"], f"{Path(attachment['file_name']).stem}.extraction.json")
        self.r2.upload(extraction_key, extraction_path, "application/json")
        self.assert_job_active(job["id"])

        # PDF retries upsert chunks in place — never delete-all first (avoids a gap).
        self.db.insert_chunks(chunks)
        self.assert_job_active(job["id"])

        has_usable_text = bool(meta.get("has_usable_text"))
        ready_at = now_iso()
        stage_errors = {}
        if meta.get("edgeparse_failed"):
            stage_errors["text"] = {
                "code": "edgeparse_failed",
                "message": meta.get("edgeparse_error") or "EdgeParse conversion failed.",
            }
            has_usable_text = False
        elif not has_usable_text:
            stage_errors["text"] = {
                "code": "no_usable_digital_text",
                "message": "No usable digital text was found in the PDF.",
            }

        meta_out = {
            "mode": "edgeparse_text",
            "progress": 100,
            "stage": "text_ready" if has_usable_text else "text_incapable",
            "page_count": meta.get("page_count", page_count),
            "word_count": meta.get("word_count", 0),
            "has_usable_text": has_usable_text,
            "extractor": "edgeparse",
            "file_name": attachment["file_name"],
            "content_type": attachment["content_type"],
            "size_bytes": attachment["size_bytes"],
            "warnings": list(stage_errors.values()),
        }
        document_patch = {
            "page_count": meta.get("page_count", page_count),
            "word_count": meta.get("word_count", 0),
            "extraction_key": extraction_key,
            "metadata": meta_out,
            "error": None,
        }
        if has_usable_text:
            document_patch["text_ready_at"] = ready_at
        if stage_errors:
            document_patch["stage_errors"] = stage_errors

        print(
            f"job {job['id']} PDF extract timings pages={page_count} "
            f"edgeparse={extract_seconds:.2f}s total={time.monotonic() - job_started:.2f}s "
            f"usable_text={has_usable_text}",
            flush=True,
        )
        return {
            "document_file_id": doc["id"],
            "status": "text_ready" if has_usable_text else "text_incapable",
            "page_count": meta.get("page_count", page_count),
            "word_count": meta.get("word_count", 0),
            "has_usable_text": has_usable_text,
            "_document_patch": document_patch,
        }

    def enrich_pdf_job(self, job, tmp):
        job_started = time.monotonic()
        doc = self.db.get_document_file(job["document_file_id"])
        attachment = self.db.get_attachment(doc["attachment_id"])
        source = tmp / safe_name(attachment["file_name"])
        self.r2.download(attachment["object_key"], source)
        self.assert_job_active(job["id"])
        limits = {**self.default_limits, **((job.get("input") or {}).get("limits") or {})}

        self.db.update_document_file(doc["id"], {
            "processing_status": "processing",
            "metadata": {
                "mode": "visual_pages",
                "progress": 2,
                "stage": "validating",
                "file_name": attachment["file_name"],
                "content_type": attachment["content_type"],
                "size_bytes": attachment["size_bytes"],
            },
        })

        source_kind = str(doc.get("kind") or "").lower()
        visual_source = source
        if source_kind in ("docx", "xlsx", "pptx"):
            self.db.update_document_file(doc["id"], {
                "metadata": {
                    "mode": "visual_pages",
                    "progress": 5,
                    "stage": "converting_to_pdf",
                    "source_kind": source_kind,
                },
            })
            visual_source = self.libreoffice_convert(source, tmp, "pdf")
            self.assert_job_active(job["id"])
        elif source_kind != "pdf":
            raise RuntimeError(f"unsupported_visual_document_kind: {source_kind}")

        reader = PdfReader(str(visual_source))
        if reader.is_encrypted:
            raise RuntimeError("password_protected")
        page_count = len(reader.pages)
        max_pages = self.limit(limits, "max_pdf_pages", 100)
        if page_count > max_pages:
            raise RuntimeError(f"too_many_pages: PDF has {page_count} pages; limit is {max_pages}")

        render_dir = tmp / "pages"
        render_dir.mkdir(parents=True, exist_ok=True)
        render_dpi = self.limit(limits, "visual_page_dpi", self.visual_page_dpi)
        self.db.update_document_file(doc["id"], {
            "metadata": {
                "mode": "visual_pages",
                "progress": 10,
                "stage": "rendering_pages",
                "page_count": page_count,
                "source_kind": source_kind,
                "embedding_model": self.embeddings.model if self.embeddings.enabled else None,
            },
        })

        ranges = split_page_ranges(page_count, self.pdf_render_workers) or [(1, page_count)]
        page_specs = []
        rendered_count = 0
        done_indexes = set()
        render_seconds = 0.0
        upload_seconds = 0.0

        def upload_and_publish(spec):
            self.r2.upload(spec["key"], spec["image_path"], "image/jpeg")
            page_row = {
                "user_id": doc["user_id"],
                "document_file_id": doc["id"],
                "page_number": spec["index"],
                "source_label": f"Page {spec['index']}",
                "image_key": spec["key"],
                "image_content_type": "image/jpeg",
                "width_px": spec["width_px"],
                "height_px": spec["height_px"],
                "text": "",
                "char_count": 0,
                "token_estimate": 0,
                "embedding": None,
                "embedding_model": None,
                "embedding_dimensions": None,
                "metadata": {
                    "page": spec["index"],
                    "source_kind": source_kind,
                    "source_etag": attachment.get("etag"),
                },
            }
            try:
                self.db.insert_pages([page_row])
            except Exception:
                self.r2.delete(spec["key"])
                raise
            return spec["index"]

        def render_range(range_index, first_page, last_page):
            range_dir = render_dir / f"range-{range_index}"
            range_dir.mkdir(parents=True, exist_ok=True)
            started = time.monotonic()
            paths = self.render_pdf_pages(
                visual_source,
                range_dir,
                render_dpi,
                page_count=(last_page - first_page) + 1,
                first_page=first_page,
                last_page=last_page,
            )
            return first_page, last_page, paths, time.monotonic() - started

        with ThreadPoolExecutor(max_workers=min(len(ranges), self.pdf_render_workers)) as render_pool:
            render_futures = [
                render_pool.submit(render_range, index, first_page, last_page)
                for index, (first_page, last_page) in enumerate(ranges, start=1)
            ]
            for render_future in as_completed(render_futures):
                first_page, last_page, image_paths, elapsed = render_future.result()
                render_seconds = max(render_seconds, elapsed)
                expected = (last_page - first_page) + 1
                if len(image_paths) != expected:
                    raise RuntimeError(
                        f"pdf_render_failed: pages {first_page}-{last_page} expected {expected}, rendered {len(image_paths)}"
                    )
                self.assert_job_active(job["id"])

                rendered_specs = []
                for offset, image_path in enumerate(image_paths):
                    page_number = first_page + offset
                    width_px, height_px = self.estimated_page_pixels(reader.pages[page_number - 1])
                    rendered_specs.append({
                        "index": page_number,
                        "image_path": image_path,
                        "key": f"users/{doc['user_id']}/documents/{doc['id']}/pages/page-{page_number:04d}.jpg",
                        "width_px": width_px,
                        "height_px": height_px,
                    })
                page_specs.extend(rendered_specs)
                rendered_count += len(rendered_specs)

                range_upload_started = time.monotonic()
                with ThreadPoolExecutor(
                    max_workers=min(self.page_upload_workers, max(1, len(rendered_specs)))
                ) as upload_pool:
                    upload_futures = [upload_pool.submit(upload_and_publish, spec) for spec in rendered_specs]
                    for upload_future in as_completed(upload_futures):
                        done_indexes.add(upload_future.result())
                upload_seconds += time.monotonic() - range_upload_started

                completed = len(done_indexes)
                self.assert_job_active(job["id"])
                if not self.db.get_document_file(doc["id"]):
                    raise RuntimeError("document_deleted")
                progress = 10 + int((completed / max(1, page_count)) * 65)
                self.db.update_document_file(doc["id"], {
                    "metadata": {
                        "mode": "visual_pages",
                        "progress": min(progress, 75),
                        "stage": "uploading_pages" if completed < page_count else "publishing_visual_manifest",
                        "page_count": page_count,
                        "source_kind": source_kind,
                        "pages_rendered": rendered_count,
                        "pages_uploaded": completed,
                        "embedding_model": self.embeddings.model if self.embeddings.enabled else None,
                    },
                })
        if len(done_indexes) < page_count:
            raise RuntimeError(f"pdf_page_upload_failed: expected {page_count} pages, uploaded {len(done_indexes)}")
        page_specs.sort(key=lambda spec: spec["index"])

        self.assert_job_active(job["id"])
        visual_ready_metadata = {
            "mode": "visual_pages",
            "progress": 75,
            "stage": "visual_ready",
            "page_count": page_count,
            "source_kind": source_kind,
            "pages_uploaded": page_count,
            "embedding_model": self.embeddings.model if self.embeddings.enabled else None,
            "file_name": attachment["file_name"],
            "content_type": attachment["content_type"],
            "size_bytes": attachment["size_bytes"],
        }
        published = self.db.publish_document_visual_ready(
            job["id"],
            self.worker_id,
            doc["id"],
            page_count,
            visual_ready_metadata,
        )
        if published is None:
            raise LeaseLostError("job_lease_lost")
        visual_ready_at = published.get("visual_ready_at") or now_iso()

        warnings = []
        stage_errors = {}
        embedding_seconds = 0.0
        embeddings = [None] * page_count
        enriched = False

        if self.embeddings.enabled:
            self.db.update_document_file(doc["id"], {
                "metadata": {
                    "mode": "visual_pages",
                    "progress": 80,
                    "stage": "embedding",
                    "page_count": page_count,
                    "source_kind": source_kind,
                    "pages_uploaded": page_count,
                    "embedding_model": self.embeddings.model,
                },
            })
            embedding_started = time.monotonic()
            try:
                embedded = self.embeddings.embed_images([spec["image_path"] for spec in page_specs])
                if not isinstance(embedded, list):
                    embedded = list(embedded or [])
                embeddings = list(embedded[:page_count])
                while len(embeddings) < page_count:
                    embeddings.append(None)
                for spec, embedding in zip(page_specs, embeddings):
                    if embedding is None:
                        continue
                    self.db.update_page(doc["id"], spec["index"], {
                        "embedding": embedding,
                        "embedding_model": self.embeddings.model,
                        "embedding_dimensions": 768,
                    })
                enriched = all(embeddings[i] is not None for i in range(page_count))
                if not enriched:
                    warnings.append({
                        "code": "embedding_incomplete",
                        "message": "One or more page embeddings were missing or skipped.",
                    })
                    stage_errors["enrichment"] = {
                        "code": "embedding_incomplete",
                        "message": "One or more page embeddings were missing or skipped.",
                    }
            except Exception as exc:
                warnings.append({"code": "embedding_failed", "message": str(exc)})
                stage_errors["enrichment"] = {
                    "code": "embedding_failed",
                    "message": str(exc),
                }
                enriched = False
            embedding_seconds = time.monotonic() - embedding_started
        else:
            warnings.append({
                "code": "embedding_unavailable",
                "message": "Image embeddings are disabled; visual pages are ready without enrichment.",
            })
            stage_errors["enrichment"] = {
                "code": "embedding_unavailable",
                "message": "Image embeddings are disabled; visual pages are ready without enrichment.",
            }
            enriched = False

        self.assert_job_active(job["id"])
        enriched_at = now_iso() if enriched else None
        meta = {
            **(doc.get("metadata") or {}),
            "mode": "visual_pages",
            "progress": 100,
            "stage": "visual_ready" if enriched else "visual_ready_embeddings_degraded",
            "page_count": page_count,
            "source_kind": source_kind,
            "pages_uploaded": page_count,
            "embedding_model": self.embeddings.model if self.embeddings.enabled else None,
            "embedding_dimensions": 768 if any(embeddings) else None,
            "warnings": warnings,
            "file_name": attachment["file_name"],
            "content_type": attachment["content_type"],
            "size_bytes": attachment["size_bytes"],
        }
        document_patch = {
            "visual_ready_at": visual_ready_at,
            "page_count": page_count,
            "metadata": meta,
            "error": None,
        }
        if enriched_at:
            document_patch["enriched_at"] = enriched_at
        if stage_errors:
            document_patch["stage_errors"] = stage_errors

        print(
            f"job {job['id']} visual enrich timings kind={source_kind} pages={page_count} "
            f"render={render_seconds:.2f}s embeddings={embedding_seconds:.2f}s "
            f"page_uploads={upload_seconds:.2f}s total={time.monotonic() - job_started:.2f}s "
            f"enriched={enriched}",
            flush=True,
        )
        return {
            "document_file_id": doc["id"],
            "status": "visual_ready",
            "page_count": page_count,
            "pages_uploaded": page_count,
            "enriched": enriched,
            "_document_patch": document_patch,
        }

    def render_page_job(self, job, tmp):
        doc = self.db.get_document_file(job["document_file_id"])
        attachment_id = doc.get("attachment_id") or (job.get("input") or {}).get("attachment_id")
        attachment = self.db.get_attachment(attachment_id)
        if not attachment:
            raise RuntimeError("attachment_not_found")
        source = tmp / safe_name(attachment["file_name"])
        self.r2.download(attachment["object_key"], source)
        self.assert_job_active(job["id"])

        source_kind = str(doc.get("kind") or "").lower()
        if source_kind in ("docx", "xlsx", "pptx"):
            visual_source = self.libreoffice_convert(source, tmp, "pdf")
            self.assert_job_active(job["id"])
        elif source_kind == "pdf":
            visual_source = source
        else:
            raise RuntimeError(f"unsupported_visual_document_kind: {source_kind}")

        input_data = job.get("input") or {}
        try:
            page_number = int(input_data.get("page_number"))
        except (TypeError, ValueError):
            raise RuntimeError("invalid_page_number")
        if page_number < 1:
            raise RuntimeError("invalid_page_number")

        reader = PdfReader(str(visual_source))
        if reader.is_encrypted:
            raise RuntimeError("password_protected")
        page_count = len(reader.pages)
        if page_number > page_count:
            raise RuntimeError("page_out_of_range")

        limits = {**self.default_limits, **(input_data.get("limits") or {})}
        render_dir = tmp / "pages"
        render_dir.mkdir(parents=True, exist_ok=True)
        render_dpi = self.limit(limits, "visual_page_dpi", self.visual_page_dpi)
        self.assert_job_active(job["id"])
        image_paths = self.render_pdf_pages(
            visual_source,
            render_dir,
            render_dpi,
            page_count=1,
            first_page=page_number,
            last_page=page_number,
        )
        if not image_paths:
            raise RuntimeError(f"pdf_render_failed: page {page_number} was not rendered")
        self.assert_job_active(job["id"])

        image_path = image_paths[0]
        key = f"users/{doc['user_id']}/documents/{doc['id']}/pages/page-{page_number:04d}.jpg"
        width_px, height_px = self.estimated_page_pixels(reader.pages[page_number - 1])
        self.r2.upload(key, image_path, "image/jpeg")
        self.assert_job_active(job["id"])
        self.db.insert_pages([{
            "user_id": doc["user_id"],
            "document_file_id": doc["id"],
            "page_number": page_number,
            "source_label": f"Page {page_number}",
            "image_key": key,
            "image_content_type": "image/jpeg",
            "width_px": width_px,
            "height_px": height_px,
            "text": "",
            "char_count": 0,
            "token_estimate": 0,
            "embedding": None,
            "embedding_model": None,
            "embedding_dimensions": None,
            "metadata": {
                "page": page_number,
                "source_kind": source_kind,
                "source_etag": attachment.get("etag"),
                "on_demand": True,
            },
        }])
        return {
            "document_file_id": doc["id"],
            "page_number": page_number,
            "image_key": key,
            "width_px": width_px,
            "height_px": height_px,
            "status": "rendered",
        }

    def limit(self, limits, key, fallback):
        try:
            value = int((limits or {}).get(key, fallback))
            return value if value > 0 else fallback
        except (TypeError, ValueError):
            return fallback

    def cap_chunks(self, chunks, limits):
        max_chars = self.limit(limits, "max_extracted_chars", self.max_extracted_chars)
        capped = []
        total = 0
        for chunk in chunks:
            text = chunk.get("text") or ""
            if total >= max_chars:
                break
            remaining = max_chars - total
            if len(text) > remaining:
                chunk = {**chunk, "text": truncate(text, remaining)}
                chunk["char_count"] = len(chunk["text"])
                chunk["token_estimate"] = max(1, len(chunk["text"]) // 4)
            capped.append(chunk)
            total += len(chunk.get("text") or "")
        return capped

    def extract(self, path, kind, user_id, document_file_id, limits=None):
        limits = {**self.default_limits, **(limits or {})}
        if kind == "pdf":
            chunks, meta = self.extract_pdf(path, user_id, document_file_id, limits)
            return self.cap_chunks(chunks, limits), meta
        if kind == "docx":
            chunks, meta = self.extract_docx(path, user_id, document_file_id, limits)
            return self.cap_chunks(chunks, limits), meta
        if kind == "xlsx":
            chunks, meta = self.extract_xlsx(path, user_id, document_file_id, limits)
            return self.cap_chunks(chunks, limits), meta
        if kind == "pptx":
            chunks, meta = self.extract_pptx(path, user_id, document_file_id, limits)
            return self.cap_chunks(chunks, limits), meta
        if kind in ("csv", "tsv"):
            chunks, meta = self.extract_csv(path, kind, user_id, document_file_id, limits)
            return self.cap_chunks(chunks, limits), meta
        raise RuntimeError(f"Unsupported document kind: {kind}")

    def chunk(self, user_id, document_file_id, index, source_type, label, text, metadata=None):
        text = truncate(text, 12000)
        return {
            "user_id": user_id,
            "document_file_id": document_file_id,
            "chunk_index": index,
            "source_type": source_type,
            "source_label": label,
            "text": text,
            "char_count": len(text),
            "token_estimate": max(1, len(text) // 4),
            "metadata": metadata or {},
        }

    def render_pdf_pages(self, path, output_dir, dpi=None, page_count=None, first_page=None, last_page=None):
        prefix = output_dir / "page"
        render_dpi = max(72, min(int(dpi or self.visual_page_dpi), 180))
        if first_page is not None and last_page is not None:
            ranges = [(int(first_page), int(last_page))]
        else:
            ranges = split_page_ranges(page_count, self.pdf_render_workers) if page_count else []

        def run_pdftoppm(first=None, last=None):
            cmd = build_pdftoppm_command(path, prefix, render_dpi, first=first, last=last)
            subprocess.run(
                cmd,
                check=True,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
            )

        if len(ranges) <= 1:
            if ranges:
                run_pdftoppm(ranges[0][0], ranges[0][1])
            else:
                run_pdftoppm()
        else:
            with ThreadPoolExecutor(max_workers=len(ranges)) as pool:
                futures = [pool.submit(run_pdftoppm, first, last) for first, last in ranges]
                for future in as_completed(futures):
                    future.result()

        def page_number(file_path):
            match = re.search(r"-(\d+)\.jpe?g$", file_path.name, re.I)
            return int(match.group(1)) if match else 0

        return sorted(output_dir.glob("page-*.jpg"), key=page_number)

    def estimated_page_pixels(self, page):
        try:
            width_pt = float(page.mediabox.width)
            height_pt = float(page.mediabox.height)
            scale = max(72, min(self.visual_page_dpi, 180)) / 72
            return max(1, round(width_pt * scale)), max(1, round(height_pt * scale))
        except Exception:
            return None, None

    def extract_pdf(self, path, user_id, document_file_id, limits, page_count=None):
        reader = PdfReader(str(path))
        if reader.is_encrypted:
            raise RuntimeError("password_protected")
        pdf_page_count = len(reader.pages)
        max_pages = self.limit(limits, "max_pdf_pages", 100)
        if pdf_page_count > max_pages:
            raise RuntimeError(f"too_many_pages: PDF has {pdf_page_count} pages; limit is {max_pages}")

        # pypdf page count is authoritative for the manifest.
        resolved_page_count = int(page_count if page_count is not None else pdf_page_count)
        if resolved_page_count < 0:
            resolved_page_count = 0

        try:
            raw = edgeparse.convert(str(path), format="json")
            grouped = group_edgeparse_pages(raw, max_page_count=resolved_page_count)
        except Exception as exc:
            # Soft text-incapable outcome: visual enrich can still succeed independently.
            return [], {
                "page_count": resolved_page_count or pdf_page_count,
                "word_count": 0,
                "has_usable_text": False,
                "extractor": "edgeparse",
                "edgeparse_failed": True,
                "edgeparse_error": str(exc),
            }

        chunks = []
        for page in grouped["pages"]:
            text = (page.get("text") or "").strip()
            if not text:
                continue
            chunks.append(self.chunk(
                user_id,
                document_file_id,
                len(chunks),
                "page",
                f"Page {page['page_number']}",
                text,
                {"page": page["page_number"], "extractor": "edgeparse"},
            ))
            if len(chunks) >= 1000:
                break

        capped = self.cap_chunks(chunks, limits)
        combined = "\n\n".join((chunk.get("text") or "").strip() for chunk in capped).strip()
        total_words = sum(len((chunk.get("text") or "").split()) for chunk in capped)
        has_usable_text = is_usable_extracted_text(combined)
        meta = {
            "page_count": resolved_page_count or pdf_page_count,
            "word_count": total_words,
            "has_usable_text": has_usable_text,
            "extractor": "edgeparse",
        }
        return capped, meta

    def extract_docx(self, path, user_id, document_file_id, limits):
        doc = Document(str(path))
        chunks = []
        buffer = []
        words = 0
        section = "Document"
        max_words = self.limit(limits, "max_docx_words", 80000)
        for paragraph in doc.paragraphs:
            text = paragraph.text.strip()
            if not text:
                continue
            words += len(text.split())
            if words > max_words:
                raise RuntimeError(f"too_many_words: DOCX exceeds {max_words} extracted words")
            style = paragraph.style.name if paragraph.style else ""
            if style.lower().startswith("heading"):
                if buffer:
                    chunks.append(self.chunk(user_id, document_file_id, len(chunks), "paragraph", section, "\n".join(buffer)))
                    buffer = []
                section = text[:120]
            else:
                buffer.append(text)
                if sum(len(x) for x in buffer) > 6000:
                    chunks.append(self.chunk(user_id, document_file_id, len(chunks), "paragraph", section, "\n".join(buffer)))
                    buffer = []
        if buffer:
            chunks.append(self.chunk(user_id, document_file_id, len(chunks), "paragraph", section, "\n".join(buffer)))
        for table_index, table in enumerate(doc.tables, start=1):
            rows = []
            for row in table.rows[:80]:
                rows.append("\t".join(cell.text.strip() for cell in row.cells))
            if rows:
                chunks.append(self.chunk(user_id, document_file_id, len(chunks), "table", f"Table {table_index}", "\n".join(rows), {"table": table_index}))
        return chunks[:1000], {"word_count": words}

    def extract_xlsx(self, path, user_id, document_file_id, limits):
        wb = load_workbook(str(path), read_only=True, data_only=False)
        values_wb = load_workbook(str(path), read_only=True, data_only=True)
        max_sheets = self.limit(limits, "max_xlsx_sheets", 25)
        max_cells = self.limit(limits, "max_xlsx_cells", 250000)
        try:
            if len(wb.worksheets) > max_sheets:
                raise RuntimeError(f"too_many_sheets: workbook has {len(wb.worksheets)} sheets; limit is {max_sheets}")

            chunks = []
            used_cells = 0
            word_count = 0
            calculation = getattr(wb, "calculation", None)
            full_recalc = bool(
                getattr(calculation, "fullCalcOnLoad", False)
                or getattr(calculation, "forceFullCalc", False)
            )
            column_block_size = 25
            chunk_chars = 5500

            for sheet in wb.worksheets:
                values_sheet = values_wb[sheet.title]
                max_row = max(0, int(sheet.max_row or 0))
                max_column = max(0, int(sheet.max_column or 0))
                if not max_row or not max_column:
                    continue
                if max_row * max_column > max_cells * 20:
                    raise RuntimeError(
                        f"too_large_used_range: sheet {sheet.title!r} declares {max_row}x{max_column} cells"
                    )

                header_row = 1
                for row_index, row in enumerate(
                    sheet.iter_rows(min_row=1, max_row=min(max_row, 20), min_col=1, max_col=max_column),
                    start=1,
                ):
                    if sum(cell.value is not None for cell in row) >= 2:
                        header_row = row_index
                        break

                for column_start in range(1, max_column + 1, column_block_size):
                    column_end = min(max_column, column_start + column_block_size - 1)
                    formula_rows = sheet.iter_rows(min_row=1, max_row=max_row, min_col=column_start, max_col=column_end)
                    cached_rows = values_sheet.iter_rows(min_row=1, max_row=max_row, min_col=column_start, max_col=column_end)
                    header_line = ""
                    buffered_lines = []
                    buffered_rows = []
                    buffered_chars = 0

                    def flush_range():
                        nonlocal buffered_lines, buffered_rows, buffered_chars
                        if not buffered_rows:
                            return
                        repeat_header = bool(header_line and header_row not in buffered_rows)
                        text_lines = ([header_line] if repeat_header else []) + buffered_lines
                        row_start = min(buffered_rows)
                        row_end = max(buffered_rows)
                        cell_range = (
                            f"{get_column_letter(column_start)}{row_start}:"
                            f"{get_column_letter(column_end)}{row_end}"
                        )
                        chunks.append(self.chunk(
                            user_id,
                            document_file_id,
                            len(chunks),
                            "sheet_range",
                            f"{sheet.title} — {cell_range}",
                            "\n".join(text_lines),
                            {
                                "sheet": sheet.title,
                                "range": cell_range,
                                "row_start": row_start,
                                "row_end": row_end,
                                "row_numbers": list(buffered_rows),
                                "column_start": column_start,
                                "column_end": column_end,
                                "header_row": header_row,
                                "header_repeated": repeat_header,
                                "sheet_state": sheet.sheet_state,
                                "extractor": "openpyxl_ranges_v2",
                            },
                        ))
                        buffered_lines = []
                        buffered_rows = []
                        buffered_chars = 0

                    for row_index, (formula_row, cached_row) in enumerate(zip(formula_rows, cached_rows), start=1):
                        rendered = []
                        nonempty = False
                        for formula_cell, cached_cell in zip(formula_row, cached_row):
                            raw = formula_cell.value
                            if raw is not None:
                                nonempty = True
                                used_cells += 1
                                if used_cells > max_cells:
                                    raise RuntimeError(f"too_many_cells: workbook exceeds {max_cells} non-empty cells")
                            if isinstance(raw, str) and raw.startswith("="):
                                cached = cached_cell.value
                                value = raw if full_recalc or cached is None else f"{raw} => {cached}"
                            else:
                                value = "" if raw is None else str(raw)
                            value = truncate(value, 2000)
                            if value:
                                word_count += len(value.split())
                            rendered.append(value)
                        line = "\t".join(rendered).rstrip()
                        if row_index == header_row:
                            header_line = line
                        if not nonempty:
                            continue
                        added_chars = len(line) + 1
                        if buffered_rows and buffered_chars + added_chars > chunk_chars:
                            flush_range()
                        buffered_lines.append(line)
                        buffered_rows.append(row_index)
                        buffered_chars += added_chars
                    flush_range()

            return chunks, {
                "sheet_count": len(wb.worksheets),
                "used_cell_count": used_cells,
                "word_count": word_count,
                "range_count": len(chunks),
                "extractor": "openpyxl_ranges_v2",
                "formula_cache_trusted": not full_recalc,
            }
        finally:
            wb.close()
            values_wb.close()

    def extract_pptx(self, path, user_id, document_file_id, limits):
        prs = Presentation(str(path))
        chunks = []
        words = 0
        for slide_index, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False) and shape.text_frame:
                    text = "\n".join(p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip())
                    if text:
                        texts.append(text)
                if getattr(shape, "has_table", False):
                    rows = []
                    for row in shape.table.rows:
                        rows.append("\t".join(cell.text.strip() for cell in row.cells))
                    if rows:
                        chunks.append(self.chunk(user_id, document_file_id, len(chunks), "table", f"Slide {slide_index} table", "\n".join(rows), {"slide": slide_index}))
            slide_text = "\n".join(texts).strip()
            if slide_text:
                words += len(slide_text.split())
                chunks.append(self.chunk(user_id, document_file_id, len(chunks), "slide", f"Slide {slide_index}", slide_text, {"slide": slide_index}))
        return chunks[:500], {"page_count": len(prs.slides), "word_count": words, "slide_count": len(prs.slides)}

    def extract_csv(self, path, kind, user_id, document_file_id, limits):
        detected = from_path(str(path)).best()
        encoding = detected.encoding if detected else "utf-8"
        delimiter = "\t" if kind == "tsv" else ","
        max_rows = self.limit(limits, "max_csv_rows", 100000)
        max_columns = self.limit(limits, "max_csv_columns", 100)
        chunks = []
        rows = []
        row_count = 0
        with open(path, "r", encoding=encoding, errors="replace", newline="") as handle:
            reader = csv.reader(handle, delimiter=delimiter)
            for row in reader:
                row_count += 1
                if row_count > max_rows:
                    raise RuntimeError(f"too_many_rows: file exceeds {max_rows} rows")
                if len(row) > max_columns:
                    raise RuntimeError(f"too_many_columns: file has {len(row)} columns; limit is {max_columns}")
                rows.append("\t".join(row))
                if len(rows) >= 200:
                    chunks.append(self.chunk(user_id, document_file_id, len(chunks), "table", f"Rows {row_count - len(rows) + 1}-{row_count}", "\n".join(rows), {"rows": [row_count - len(rows) + 1, row_count]}))
                    rows = []
        if rows:
            chunks.append(self.chunk(user_id, document_file_id, len(chunks), "table", f"Rows {row_count - len(rows) + 1}-{row_count}", "\n".join(rows), {"rows": [row_count - len(rows) + 1, row_count]}))
        return chunks[:1000], {"row_count": row_count}

    def create_job(self, job, tmp):
        input_data = job.get("input") or {}
        fmt = input_data.get("format") or job["job_type"].split(".")[-1]
        title = input_data.get("title") or "Generated document"
        if fmt == "docx":
            path = self.create_js_artifact(tmp, title, input_data, "docx") or self.create_docx(tmp, title, input_data)
            content_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        elif fmt == "xlsx":
            path = self.create_xlsx(tmp, title, input_data)
            content_type = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        elif fmt == "pdf":
            path = self.create_pdf(tmp, title, input_data)
            content_type = "application/pdf"
        elif fmt == "pptx":
            path = self.create_js_artifact(tmp, title, input_data, "pptx") or self.create_pptx(tmp, title, input_data)
            content_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        else:
            raise RuntimeError(f"Unsupported create format: {fmt}")
        return self.store_generated(job, tmp, path, fmt, content_type, "generated", None)

    def create_js_artifact(self, tmp, title, input_data, fmt):
        if not USE_JS_ARTIFACT_GENERATOR:
            return None
        generator = Path(NODE_ARTIFACT_GENERATOR)
        if not generator.exists():
            return None
        payload = dict(input_data or {})
        payload["format"] = fmt
        payload["title"] = title
        input_path = tmp / f"artifact-input-{uuid.uuid4()}.json"
        outdir = tmp / f"artifact-out-{uuid.uuid4()}"
        outdir.mkdir(exist_ok=True)
        input_path.write_text(json.dumps(payload), encoding="utf-8")
        cmd = [NODE_BIN, str(generator), str(input_path), str(outdir)]
        try:
            result = subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True, timeout=90)
            data = json.loads(result.stdout or "{}")
            output = Path(data.get("path") or "")
            if not output.exists():
                raise RuntimeError("JS artifact generator did not produce an output file.")
            resolved_output = output.resolve()
            resolved_outdir = outdir.resolve()
            if resolved_outdir not in resolved_output.parents:
                raise RuntimeError("JS artifact generator returned a path outside the output directory.")
            if output.suffix.lower() != f".{fmt}":
                raise RuntimeError("JS artifact generator returned the wrong file type.")
            return output
        except Exception as exc:
            print(f"JS artifact generator failed for {fmt}; using Python fallback: {exc}", flush=True)
            return None

    def create_docx(self, tmp, title, input_data):
        doc = Document()
        section = doc.sections[0]
        section.top_margin = DocxInches(0.7)
        section.bottom_margin = DocxInches(0.7)
        section.left_margin = DocxInches(0.8)
        section.right_margin = DocxInches(0.8)
        doc.add_heading(title, level=1)
        body = strip_duplicate_title_heading(artifact_content(input_data), title)
        if body:
            self.append_docx_markdown(doc, body)
        else:
            instructions = input_data.get("instructions") or ""
            for block in instructions.split("\n\n"):
                if block.strip():
                    doc.add_paragraph(clean_markdown(block))
        for section in input_data.get("sections") or []:
            heading = section.get("heading") or section.get("title")
            if heading:
                doc.add_heading(str(heading), level=2)
            content = section.get("content") or section.get("text") or ""
            if content:
                self.append_docx_markdown(doc, str(content))
        for table in input_data.get("tables") or []:
            self.append_docx_table(doc, table)
        path = tmp / f"{safe_name(title)}.docx"
        doc.save(path)
        return path

    def append_docx_markdown(self, doc, text):
        in_code = False
        code_lines = []
        lines = str(text or "").splitlines()
        index = 0
        while index < len(lines):
            raw_line = lines[index]
            line = raw_line.strip()
            if line.startswith("```"):
                if in_code:
                    self.append_docx_code_block(doc, "\n".join(code_lines))
                    code_lines = []
                    in_code = False
                else:
                    in_code = True
                    code_lines = []
                index += 1
                continue
            if in_code:
                code_lines.append(raw_line.rstrip())
                index += 1
                continue
            if not line:
                index += 1
                continue
            table_data, next_index = collect_markdown_table(lines, index)
            if table_data:
                self.append_docx_table(doc, table_data)
                index = next_index
                continue
            if re.match(r"^-{3,}$", line):
                index += 1
                continue
            heading = re.match(r"^(#{1,3})\s+(.+)$", line)
            if heading:
                doc.add_heading(clean_markdown(heading.group(2)), level=min(len(heading.group(1)), 3))
                index += 1
                continue
            bullet = re.match(r"^[-*]\s+(.+)$", line)
            if bullet:
                doc.add_paragraph(clean_markdown(bullet.group(1)), style="List Bullet")
                index += 1
                continue
            numbered = re.match(r"^\d+[.)]\s+(.+)$", line)
            if numbered:
                doc.add_paragraph(clean_markdown(numbered.group(1)), style="List Number")
                index += 1
                continue
            doc.add_paragraph(clean_markdown(line))
            index += 1
        if in_code and code_lines:
            self.append_docx_code_block(doc, "\n".join(code_lines))

    def append_docx_code_block(self, doc, text):
        paragraph = doc.add_paragraph()
        run = paragraph.add_run(normalize_math_symbols(text).strip())
        run.font.name = "Courier New"

    def append_docx_table(self, doc, table_data):
        title = table_data.get("title") or table_data.get("caption")
        if title:
            doc.add_heading(str(title), level=2)
        rows = table_data.get("rows") or table_data.get("data") or []
        headers = table_data.get("headers") or []
        if headers:
            rows = [headers] + rows
        if not rows:
            return
        width = max(len(row) if isinstance(row, list) else 1 for row in rows)
        table = doc.add_table(rows=0, cols=max(1, width))
        table.style = "Table Grid"
        for row in rows[:200]:
            cells = table.add_row().cells
            values = normalize_table_row_width(row if isinstance(row, list) else [row], width)
            for i, value in enumerate(values[:width]):
                cells[i].text = clean_markdown(str(value))

    def create_xlsx(self, tmp, title, input_data):
        path = tmp / f"{safe_name(title)}.xlsx"
        payload = dict(input_data or {})
        payload["title"] = title
        return Path(create_xlsx_workbook(path, payload))

    def create_pptx(self, tmp, title, input_data):
        prs = Presentation()
        prs.slide_width = PptxInches(13.333)
        prs.slide_height = PptxInches(7.5)
        slides = self.presentation_slides(title, input_data)
        for index, slide_data in enumerate(slides[:40]):
            self.append_pptx_slide(prs, slide_data, is_first=(index == 0))
        path = tmp / f"{safe_name(title)}.pptx"
        prs.save(path)
        return path

    def presentation_slides(self, title, input_data):
        data = input_data.get("data") if isinstance(input_data.get("data"), dict) else {}
        slides = data.get("slides") if isinstance(data.get("slides"), list) else []
        if slides:
            return [self.normalize_slide_data(slide, title if index == 0 else "") for index, slide in enumerate(slides)]

        out = []
        content = strip_duplicate_title_heading(artifact_content(input_data), title)
        if content:
            out = self.slides_from_markdown(title, content)
        if not out:
            sections = input_data.get("sections") if isinstance(input_data.get("sections"), list) else []
            if sections:
                out.append({"title": title, "subtitle": input_data.get("instructions") or "", "bullets": []})
                for section in sections:
                    out.append({
                        "title": section.get("heading") or section.get("title") or "Section",
                        "subtitle": section.get("message") or "",
                        "bullets": self.text_to_bullets(section.get("content") or section.get("text") or "")
                    })
        if not out:
            out = [
                {"title": title, "subtitle": input_data.get("instructions") or "", "bullets": []},
                {"title": "Overview", "bullets": self.text_to_bullets(input_data.get("instructions") or "Generated presentation")}
            ]
        return out

    def normalize_slide_data(self, slide, fallback_title=""):
        if not isinstance(slide, dict):
            return {"title": clean_markdown(str(slide)), "bullets": []}
        bullets = slide.get("bullets") or slide.get("points") or slide.get("items") or []
        if isinstance(bullets, str):
            bullets = self.text_to_bullets(bullets)
        elif not isinstance(bullets, list):
            bullets = []
        return {
            "title": clean_markdown(slide.get("title") or slide.get("heading") or fallback_title or "Slide"),
            "subtitle": clean_markdown(slide.get("subtitle") or slide.get("message") or slide.get("takeaway") or ""),
            "bullets": [clean_markdown(item) for item in bullets if clean_markdown(item)][:8],
            "notes": str(slide.get("notes") or slide.get("speaker_notes") or "").strip(),
            "table": slide.get("table") if isinstance(slide.get("table"), dict) else None,
        }

    def slides_from_markdown(self, title, content):
        slides = [{"title": title, "subtitle": "", "bullets": []}]
        current = None
        for raw_line in str(content or "").splitlines():
            line = raw_line.strip()
            if not line:
                continue
            heading = re.match(r"^#{1,3}\s+(.+)$", line)
            if heading:
                current = {"title": clean_markdown(heading.group(1)), "bullets": []}
                slides.append(current)
                continue
            bullet = re.match(r"^[-*]\s+(.+)$", line) or re.match(r"^\d+[.)]\s+(.+)$", line)
            if bullet:
                if not current:
                    current = {"title": "Key Points", "bullets": []}
                    slides.append(current)
                current.setdefault("bullets", []).append(clean_markdown(bullet.group(1)))
                continue
            if not current:
                slides[0]["subtitle"] = (slides[0].get("subtitle") or line)[:180]
            elif len(current.get("bullets", [])) < 6:
                current.setdefault("bullets", []).append(clean_markdown(line))
        return [self.normalize_slide_data(slide) for slide in slides if slide.get("title") or slide.get("subtitle") or slide.get("bullets")]

    def text_to_bullets(self, text):
        bullets = []
        for line in str(text or "").splitlines():
            cleaned = clean_markdown(re.sub(r"^[-*\d.)\s]+", "", line).strip())
            if cleaned:
                bullets.append(cleaned)
        if not bullets and str(text or "").strip():
            parts = re.split(r"(?<=[.!?])\s+", clean_markdown(text))
            bullets = [part for part in parts if part][:6]
        return bullets[:8]

    def append_pptx_slide(self, prs, slide_data, is_first=False):
        slide_data = self.normalize_slide_data(slide_data)
        if is_first:
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = slide_data["title"] or "Presentation"
            subtitle = slide.placeholders[1]
            subtitle.text = slide_data.get("subtitle") or ""
            self.style_pptx_title(slide.shapes.title, size=34)
            self.style_pptx_text(subtitle, size=17, color=RGBColor(90, 90, 90))
        else:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            self.add_pptx_title(slide, slide_data["title"])
            if slide_data.get("subtitle"):
                self.add_pptx_textbox(slide, slide_data["subtitle"], 0.7, 1.15, 11.8, 0.55, size=15, color=RGBColor(80, 80, 80))
            if slide_data.get("table"):
                self.add_pptx_table(slide, slide_data["table"], 0.7, 1.9, 11.8, 4.5)
            else:
                self.add_pptx_bullets(slide, slide_data.get("bullets") or [], 1.0, 1.85, 11.0, 4.8)
        notes = slide_data.get("notes")
        if notes:
            try:
                slide.notes_slide.notes_text_frame.text = notes[:2000]
            except Exception:
                pass

    def add_pptx_title(self, slide, title):
        box = slide.shapes.add_textbox(PptxInches(0.6), PptxInches(0.35), PptxInches(12.0), PptxInches(0.65))
        frame = box.text_frame
        frame.clear()
        paragraph = frame.paragraphs[0]
        paragraph.text = clean_markdown(title or "Slide")
        paragraph.font.size = Pt(26)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(24, 24, 24)
        paragraph.alignment = PP_ALIGN.LEFT

    def add_pptx_textbox(self, slide, text, x, y, w, h, size=16, color=None):
        box = slide.shapes.add_textbox(PptxInches(x), PptxInches(y), PptxInches(w), PptxInches(h))
        box.text_frame.word_wrap = True
        box.text_frame.text = clean_markdown(text)
        self.style_pptx_text(box, size=size, color=color or RGBColor(40, 40, 40))
        return box

    def add_pptx_bullets(self, slide, bullets, x, y, w, h):
        box = slide.shapes.add_textbox(PptxInches(x), PptxInches(y), PptxInches(w), PptxInches(h))
        frame = box.text_frame
        frame.word_wrap = True
        frame.clear()
        items = bullets or ["Add the key point for this slide."]
        for index, item in enumerate(items[:8]):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = clean_markdown(item)
            paragraph.level = 0
            paragraph.font.size = Pt(18 if len(items) <= 5 else 15)
            paragraph.font.color.rgb = RGBColor(35, 35, 35)
            paragraph.space_after = Pt(8)

    def add_pptx_table(self, slide, table_data, x, y, w, h):
        rows = table_data.get("rows") or table_data.get("data") or []
        headers = table_data.get("headers") or []
        if headers:
            rows = [headers] + rows
        rows = [row if isinstance(row, list) else [row] for row in rows[:12]]
        if not rows:
            return
        cols = max(1, max(len(row) for row in rows))
        shape = slide.shapes.add_table(len(rows), cols, PptxInches(x), PptxInches(y), PptxInches(w), PptxInches(h))
        table = shape.table
        for row_index, row in enumerate(rows):
            values = normalize_table_row_width(row, cols)
            for col_index, value in enumerate(values[:cols]):
                cell = table.cell(row_index, col_index)
                cell.text = clean_markdown(str(value))
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(10 if cols > 4 else 12)
                    paragraph.font.bold = row_index == 0
                    paragraph.font.color.rgb = RGBColor(24, 24, 24)

    def style_pptx_title(self, shape, size=30):
        for paragraph in shape.text_frame.paragraphs:
            paragraph.font.size = Pt(size)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(24, 24, 24)

    def style_pptx_text(self, shape, size=16, color=None):
        for paragraph in shape.text_frame.paragraphs:
            paragraph.font.size = Pt(size)
            paragraph.font.color.rgb = color or RGBColor(45, 45, 45)

    def create_pdf(self, tmp, title, input_data):
        path = tmp / f"{safe_name(title)}.pdf"
        fonts = register_pdf_fonts()
        doc = SimpleDocTemplate(
            str(path),
            pagesize=letter,
            leftMargin=54,
            rightMargin=54,
            topMargin=54,
            bottomMargin=54,
        )
        styles = getSampleStyleSheet()
        styles["Title"].fontName = fonts["bold"]
        styles["Title"].fontSize = 20
        styles["Title"].leading = 24
        styles["Title"].spaceAfter = 12
        styles["Heading1"].fontName = fonts["bold"]
        styles["Heading1"].fontSize = 15
        styles["Heading1"].leading = 19
        styles["Heading1"].spaceBefore = 10
        styles["Heading1"].spaceAfter = 6
        styles["Heading2"].fontName = fonts["bold"]
        styles["Heading2"].fontSize = 13
        styles["Heading2"].leading = 16
        styles["Heading2"].spaceBefore = 8
        styles["Heading2"].spaceAfter = 5
        styles["Normal"].fontName = fonts["regular"]
        styles["Normal"].fontSize = 10.5
        styles["Normal"].leading = 14
        callout_style = ParagraphStyle(
            "KluiCallout",
            parent=styles["Normal"],
            fontName=fonts["regular"],
            fontSize=10.5,
            leading=14,
            leftIndent=10,
            rightIndent=10,
            spaceBefore=4,
            spaceAfter=10,
            borderColor=colors.HexColor("#CBD5E1"),
            borderWidth=0.5,
            borderPadding=8,
            backColor=colors.HexColor("#F8FAFC"),
        )
        story = [Paragraph(title, styles["Title"]), Spacer(1, 12)]
        summary = (input_data.get("data") or {}).get("summary") or (input_data.get("data") or {}).get("recommendation") or input_data.get("recommendation")
        if summary:
            story.append(Paragraph(f"<b>Key takeaway:</b> {xml_escape(clean_markdown(str(summary)))}", callout_style))
            story.append(Spacer(1, 8))
        body = strip_duplicate_title_heading(artifact_content(input_data), title)
        if body:
            self.append_pdf_markdown(story, body, styles, fonts)
        else:
            for block in (input_data.get("instructions") or "").split("\n\n"):
                if block.strip():
                    story.append(Paragraph(xml_escape(clean_markdown(block)).replace("\n", "<br/>"), styles["Normal"]))
                    story.append(Spacer(1, 8))
        for section in input_data.get("sections") or []:
            heading = section.get("heading") or section.get("title")
            content = section.get("content") or section.get("text") or ""
            if heading:
                story.append(Paragraph(xml_escape(str(heading)), styles["Heading2"]))
                story.append(Spacer(1, 6))
            if content:
                self.append_pdf_markdown(story, str(content), styles, fonts)
        for table in input_data.get("tables") or []:
            self.append_pdf_table(story, table, styles, fonts)
        def draw_footer(canvas, document):
            canvas.saveState()
            canvas.setFont(fonts["regular"], 8)
            canvas.setFillColor(colors.grey)
            canvas.drawRightString(letter[0] - 54, 24, f"Page {document.page}")
            canvas.restoreState()

        doc.build(story, onFirstPage=draw_footer, onLaterPages=draw_footer)
        return path

    def append_pdf_markdown(self, story, text, styles, fonts):
        code_style = ParagraphStyle(
            "KluiCode",
            parent=styles["Code"],
            fontName=fonts["mono"],
            fontSize=9,
            leading=12,
            backColor=colors.whitesmoke,
            borderColor=colors.lightgrey,
            borderWidth=0.25,
            borderPadding=6,
            leftIndent=8,
            rightIndent=8,
            spaceBefore=4,
            spaceAfter=8,
        )
        in_code = False
        code_lines = []
        lines = str(text or "").splitlines()
        index = 0
        while index < len(lines):
            raw_line = lines[index]
            line = raw_line.strip()
            if line.startswith("```"):
                if in_code:
                    story.append(Preformatted(xml_escape(normalize_math_symbols("\n".join(code_lines)).strip()), code_style))
                    story.append(Spacer(1, 6))
                    code_lines = []
                    in_code = False
                else:
                    in_code = True
                    code_lines = []
                index += 1
                continue
            if in_code:
                code_lines.append(raw_line.rstrip())
                index += 1
                continue
            if not line:
                story.append(Spacer(1, 6))
                index += 1
                continue
            table_data, next_index = collect_markdown_table(lines, index)
            if table_data:
                self.append_pdf_table(story, table_data, styles, fonts)
                index = next_index
                continue
            if re.match(r"^-{3,}$", line):
                story.append(Spacer(1, 8))
                index += 1
                continue
            heading = re.match(r"^(#{1,3})\s+(.+)$", line)
            if heading:
                style = styles["Heading1"] if len(heading.group(1)) == 1 else styles["Heading2"]
                story.append(Paragraph(xml_escape(clean_markdown(heading.group(2))), style))
                story.append(Spacer(1, 6))
                index += 1
                continue
            bullet = re.match(r"^[-*]\s+(.+)$", line)
            if bullet:
                story.append(Paragraph(f"- {xml_escape(clean_markdown(bullet.group(1)))}", styles["Normal"]))
                story.append(Spacer(1, 4))
                index += 1
                continue
            numbered = re.match(r"^(\d+[.)])\s+(.+)$", line)
            if numbered:
                story.append(Paragraph(f"{xml_escape(numbered.group(1))} {xml_escape(clean_markdown(numbered.group(2)))}", styles["Normal"]))
                story.append(Spacer(1, 4))
                index += 1
                continue
            story.append(Paragraph(xml_escape(clean_markdown(line)), styles["Normal"]))
            story.append(Spacer(1, 6))
            index += 1
        if in_code and code_lines:
            story.append(Preformatted(xml_escape(normalize_math_symbols("\n".join(code_lines)).strip()), code_style))
            story.append(Spacer(1, 6))

    def append_pdf_table(self, story, table_data, styles, fonts):
        title = table_data.get("title") or table_data.get("caption")
        if title:
            story.append(Paragraph(xml_escape(str(title)), styles["Heading2"]))
            story.append(Spacer(1, 6))
        rows = table_data.get("rows") or table_data.get("data") or []
        headers = table_data.get("headers") or []
        if headers:
            rows = [headers] + rows
        if not rows:
            return
        width = max(len(row) if isinstance(row, list) else 1 for row in rows)
        table_style = ParagraphStyle(
            "KluiTableCell",
            parent=styles["Normal"],
            fontName=fonts["regular"],
            fontSize=8.5,
            leading=11,
            wordWrap="CJK",
        )
        clean_rows = [
            [
                Paragraph(xml_escape(clean_markdown(str(cell))).replace("\n", "<br/>"), table_style)
                for cell in normalize_table_row_width(row if isinstance(row, list) else [row], width)
            ]
            for row in rows[:50]
        ]
        col_widths = self.pdf_table_col_widths(rows[:50], width)
        table = Table(clean_rows, colWidths=col_widths, repeatRows=1, hAlign="LEFT", splitByRow=1)
        table.setStyle(TableStyle([
            ("GRID", (0, 0), (-1, -1), 0.25, colors.grey),
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#F1F5F9")),
            ("FONTNAME", (0, 0), (-1, 0), fonts["bold"]),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#FAFAFA")]),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("LEFTPADDING", (0, 0), (-1, -1), 5),
            ("RIGHTPADDING", (0, 0), (-1, -1), 5),
            ("TOPPADDING", (0, 0), (-1, -1), 4),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
        ]))
        story.append(table)
        story.append(Spacer(1, 10))

    def pdf_table_col_widths(self, rows, width):
        if width <= 0:
            return None
        max_width = letter[0] - 108
        samples = []
        for col in range(width):
            lengths = []
            for row in rows:
                values = normalize_table_row_width(row if isinstance(row, list) else [row], width)
                lengths.append(min(len(str(values[col] or "")), 80))
            samples.append(max([8] + lengths))
        total = sum(samples) or width
        raw = [max_width * (sample / total) for sample in samples]
        min_width = min(72, max_width / width)
        widths = [max(min_width, value) for value in raw]
        scale = max_width / sum(widths)
        return [value * scale for value in widths]

    def edit_job(self, job, tmp):
        source_doc = self.db.get_document_file(job["document_file_id"])
        attachment = self.db.get_attachment(source_doc["attachment_id"])
        source = tmp / safe_name(attachment["file_name"])
        self.r2.download(attachment["object_key"], source)
        kind = source_doc["kind"]
        if kind == "docx":
            output = self.edit_docx(source, tmp, job.get("input") or {})
            content_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        elif kind == "xlsx":
            output = self.edit_xlsx(source, tmp, job.get("input") or {})
            content_type = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        else:
            raise RuntimeError("Editing this document type is not supported yet.")
        return self.store_generated(job, tmp, output, kind, content_type, "edited", source_doc)

    def edit_docx(self, source, tmp, input_data):
        doc = Document(str(source))
        operations = input_data.get("operations") or []
        changed = False
        for op in operations:
            find = str(op.get("find") or op.get("old_text") or "")
            replace = str(op.get("replace") or op.get("new_text") or "")
            if not find:
                continue
            for paragraph in doc.paragraphs:
                if find in paragraph.text:
                    paragraph.text = paragraph.text.replace(find, replace)
                    changed = True
        instructions = input_data.get("instructions") or ""
        if instructions and not changed:
            doc.add_page_break()
            doc.add_heading("Requested edits", level=1)
            doc.add_paragraph(instructions)
        output = tmp / f"edited-{source.name}"
        doc.save(output)
        return output

    def edit_xlsx(self, source, tmp, input_data):
        wb = load_workbook(str(source))
        operations = input_data.get("operations") or []
        if not operations:
            raise RuntimeError("xlsx_edit_requires_operations")
        if len(operations) > 100:
            raise RuntimeError("xlsx_edit_too_many_operations")
        changed_cells = 0

        def require_sheet(op):
            sheet_name = str(op.get("sheet") or "").strip()
            if not sheet_name or sheet_name not in wb.sheetnames:
                raise RuntimeError(f"xlsx_sheet_not_found: {sheet_name or '(missing)'}")
            return wb[sheet_name]

        def bounds(op, key="range"):
            reference = str(op.get(key) or "").strip().upper()
            if not reference:
                raise RuntimeError(f"xlsx_edit_missing_{key}")
            try:
                return reference, range_boundaries(reference)
            except ValueError as exc:
                raise RuntimeError(f"xlsx_edit_invalid_{key}: {reference}") from exc

        def safe_value(value):
            return validate_formula(value) if isinstance(value, str) and value.startswith("=") else value

        for op in operations:
            if not isinstance(op, dict):
                raise RuntimeError("xlsx_edit_operation_must_be_object")
            operation = str(op.get("type") or "").strip().lower()

            if operation == "add_sheet":
                name = str(op.get("name") or "").strip()
                if not name or len(name) > 31 or name in wb.sheetnames:
                    raise RuntimeError(f"xlsx_invalid_new_sheet: {name or '(missing)'}")
                wb.create_sheet(name)
                continue

            ws = require_sheet(op)
            if operation == "rename_sheet":
                name = str(op.get("new_name") or "").strip()
                if not name or len(name) > 31 or name in wb.sheetnames:
                    raise RuntimeError(f"xlsx_invalid_new_sheet: {name or '(missing)'}")
                ws.title = name
            elif operation == "delete_sheet":
                if len(wb.sheetnames) == 1:
                    raise RuntimeError("xlsx_cannot_delete_last_sheet")
                wb.remove(ws)
            elif operation in ("set_cell", "set_formula"):
                reference, (min_col, min_row, max_col, max_row) = bounds(op, "cell")
                if min_col != max_col or min_row != max_row:
                    raise RuntimeError(f"xlsx_edit_cell_must_be_single: {reference}")
                if operation == "set_cell" and "value" not in op:
                    raise RuntimeError("xlsx_edit_missing_value")
                value = op.get("formula") if operation == "set_formula" else op.get("value")
                ws.cell(min_row, min_col).value = safe_value(value)
                changed_cells += 1
            elif operation == "set_range":
                reference, (min_col, min_row, max_col, max_row) = bounds(op)
                values = op.get("values")
                expected_rows = max_row - min_row + 1
                expected_columns = max_col - min_col + 1
                if not isinstance(values, list) or len(values) != expected_rows:
                    raise RuntimeError(f"xlsx_edit_range_row_mismatch: {reference}")
                for row_offset, row in enumerate(values):
                    if not isinstance(row, list) or len(row) != expected_columns:
                        raise RuntimeError(f"xlsx_edit_range_column_mismatch: {reference}")
                    for column_offset, value in enumerate(row):
                        ws.cell(min_row + row_offset, min_col + column_offset).value = safe_value(value)
                        changed_cells += 1
            elif operation == "append_rows":
                rows = op.get("rows")
                if not rows or not isinstance(rows, list) or not all(isinstance(row, list) for row in rows):
                    raise RuntimeError("xlsx_edit_rows_must_be_arrays")
                for row in rows:
                    ws.append([safe_value(value) for value in row])
                    changed_cells += len(row)
            elif operation == "clear_range":
                reference, (min_col, min_row, max_col, max_row) = bounds(op)
                for row in ws.iter_rows(min_row=min_row, max_row=max_row, min_col=min_col, max_col=max_col):
                    for cell in row:
                        cell.value = None
                        changed_cells += 1
            elif operation == "set_number_format":
                reference, (min_col, min_row, max_col, max_row) = bounds(op)
                number_format = str(op.get("format") or "").strip()
                if not number_format or len(number_format) > 100:
                    raise RuntimeError("xlsx_edit_invalid_number_format")
                for row in ws.iter_rows(min_row=min_row, max_row=max_row, min_col=min_col, max_col=max_col):
                    for cell in row:
                        cell.number_format = number_format
                        changed_cells += 1
            else:
                raise RuntimeError(f"xlsx_edit_unknown_operation: {operation or '(missing)'}")

            if changed_cells > 250000:
                raise RuntimeError("xlsx_edit_too_many_cells")

        output = tmp / f"edited-{source.name}"
        wb.calculation.fullCalcOnLoad = True
        wb.calculation.forceFullCalc = True
        wb.calculation.calcMode = "auto"
        wb.save(output)
        wb.close()
        check = load_workbook(str(output), read_only=True, data_only=False)
        check.close()
        return output

    def export_job(self, job, tmp):
        source_doc = self.db.get_document_file(job["document_file_id"])
        attachment = self.db.get_attachment(source_doc["attachment_id"])
        source = tmp / safe_name(attachment["file_name"])
        self.r2.download(attachment["object_key"], source)
        target = (job.get("input") or {}).get("target_format") or "pdf"
        if source_doc["kind"] == target:
            output = tmp / source.name
            shutil.copyfile(source, output)
        elif target == "pdf" and source_doc["kind"] in ("docx", "xlsx", "pptx"):
            output = self.libreoffice_convert(source, tmp, "pdf")
        else:
            raise RuntimeError("Unsupported export conversion.")
        content_type = "application/pdf" if target == "pdf" else attachment["content_type"]
        return self.store_generated(job, tmp, output, target, content_type, "exported", source_doc)

    def libreoffice_convert(self, source, tmp, target):
        outdir = tmp / "out"
        outdir.mkdir(exist_ok=True)
        profile = tmp / f"lo-{uuid.uuid4()}"
        cmd = [
            "soffice",
            "--headless",
            f"-env:UserInstallation=file://{profile}",
            "--convert-to",
            target,
            "--outdir",
            str(outdir),
            str(source),
        ]
        subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=90)
        converted = outdir / f"{source.stem}.{target}"
        if not converted.exists():
            raise RuntimeError("LibreOffice conversion failed.")
        return converted

    def store_generated(self, job, tmp, path, kind, content_type, source, parent_doc):
        user_id = job["user_id"]
        input_data = job.get("input") or {}
        preview = bool(input_data.get("preview"))
        editor_markdown = str(input_data.get("editor_markdown") or "").strip()
        editor_metadata = {
            "editor_markdown": editor_markdown,
            "editor_revision": 1,
            "editable": True,
        } if editor_markdown and kind in ("docx", "pdf") and not preview else {}
        key = self.object_key(user_id, path.name)
        etag = self.r2.upload(key, path, content_type)
        attachment = self.db.create_attachment({
            "user_id": user_id,
            "conversation_id": job.get("conversation_id"),
            "message_id": job.get("message_id"),
            "category": "document",
            "object_key": key,
            "file_name": path.name,
            "content_type": content_type,
            "size_bytes": path.stat().st_size,
            "etag": etag,
            "status": "uploaded",
            "uploaded_at": now_iso(),
        })
        document_file = self.db.create_document_file({
            "attachment_id": attachment["id"],
            "user_id": user_id,
            "conversation_id": job.get("conversation_id"),
            "message_id": job.get("message_id"),
            "kind": kind,
            "source": source,
            "parent_document_id": parent_doc["id"] if parent_doc else None,
            "version_no": int(parent_doc.get("version_no", 0)) + 1 if parent_doc else 1,
            "source_etag": etag,
            "processing_status": "processing",
            "metadata": {"generated_by_job": job["id"], **editor_metadata, **({"preview": True} if preview else {})},
        })
        chunks, meta = self.extract(path, kind, user_id, document_file["id"])
        self.db.insert_chunks(chunks)
        ready_at = now_iso()
        # Generated artifacts are created outside the upload extract/enrich pair, so the
        # worker still marks text readiness (and legacy ready) on the new document row.
        self.db.update_document_file(document_file["id"], {
            "processing_status": "ready",
            "text_ready_at": ready_at,
            "page_count": meta.get("page_count"),
            "word_count": meta.get("word_count"),
            "sheet_count": meta.get("sheet_count"),
            "used_cell_count": meta.get("used_cell_count"),
            "metadata": {**meta, "generated_by_job": job["id"], **editor_metadata, **({"preview": True} if preview else {})},
            "error": None,
        })
        return {
            "attachment_id": attachment["id"],
            "document_file_id": document_file["id"],
            "file_name": attachment["file_name"],
            "kind": kind,
            "status": "ready",
            **({"preview": True} if preview else {}),
            "download_url": f"/api/attachments/{attachment['id']}/download",
        }


def worker_concurrency():
    return env_int("DOCUMENT_WORKER_CONCURRENCY", 1, minimum=1, maximum=WORKER_CONCURRENCY_CAP)


def main():
    concurrency = worker_concurrency()
    if concurrency <= 1:
        Processor().run()
        return

    print(f"starting {concurrency} document worker loops", flush=True)
    threads = []
    for _ in range(concurrency):
        processor = Processor()
        thread = threading.Thread(target=processor.run, name=processor.worker_id, daemon=True)
        thread.start()
        threads.append(thread)
    for thread in threads:
        thread.join()


if __name__ == "__main__":
    main()
